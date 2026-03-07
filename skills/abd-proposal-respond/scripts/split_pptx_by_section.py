#!/usr/bin/env python3
"""
Split a PowerPoint file into separate files by INNER section number.

Splits by slide titles like "SECTION 1.1", "SECTION 2.6", "SECTION 4.1" (inner)
NOT by outer section numbers (1, 2, 3, 4).

Output files are saved in the same directory as the source file.

Usage:
  python split_pptx_by_section.py <path_to.pptx>
  python split_pptx_by_section.py <path_to.pptx> --output-dir <dir>

Requires: pywin32, python-pptx (pip install pywin32 python-pptx)
"""

import os
import re
import sys
import tempfile
import time
from pathlib import Path

try:
    import win32com.client
except ImportError:
    win32com = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def _sanitize_filename(name: str) -> str:
    """Make section name safe for use as filename (e.g. 'The Agile Mindset: Achieving Flow' -> 'The Agile Mindset - Achieving Flow')."""
    safe = re.sub(r'[<>:"/\\|?*]', " - ", name)
    safe = re.sub(r"\s+", " ", safe).strip()
    safe = safe.strip(". ")
    return safe or "Section"


# Match "SECTION 1.1", "SECTION 2.6", "Section 4.1", "SECTION 3" (no decimal)
_SECTION_RE = re.compile(
    r"\b(?:SECTION|Section)\s+(\d+(?:\.\d+)?)\b",
    re.IGNORECASE,
)

# Titles to NEVER use as section names (master/header text that repeats on every slide)
_BAD_TITLES = frozenset(
    s.lower()
    for s in (
        "EY AGILE CURRICULUM",
        "EY AGILE",
    )
)


def _is_bad_title(text: str) -> bool:
    """True if this is header/master text we should never use as a section name."""
    if not text or not text.strip():
        return True
    t = text.strip()
    if t.lower() in _BAD_TITLES:
        return True
    if re.match(r"^(?:SECTION|Section)\s+\d+(?:\.\d+)?\s*$", t, re.I):
        return True
    return False


def _get_all_slide_text(slide) -> str:
    """Concatenate all text from slide shapes (for SECTION pattern detection)."""
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text and sh.text.strip():
            parts.append(sh.text.strip())
    return "\n".join(parts)


def _get_slide_title(slide, prs=None) -> str:
    """Extract title from the CENTER of the slide (e.g. 'The Agile Mindset: Achieving Flow').
    Never use header/footer or top-left master text like 'EY AGILE CURRICULUM'."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        PP_PLACEHOLDER = None

    def _skip_placeholder(sh) -> bool:
        """Exclude header, footer, date, slide number - these repeat on every slide."""
        if not getattr(sh, "is_placeholder", False):
            return False
        pf = getattr(sh, "placeholder_format", None)
        if not pf:
            return False
        ptype = getattr(pf, "type", None)
        if PP_PLACEHOLDER and ptype is not None:
            return ptype in (
                PP_PLACEHOLDER.HEADER,
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.DATE,
                PP_PLACEHOLDER.SLIDE_NUMBER,
            )
        return False

    # Use actual slide dimensions if available, else standard 10" x 7.5"
    if prs:
        slide_height = getattr(prs, "slide_height", None) or 6858000
    else:
        slide_height = 6858000

    def _center_y(sh):
        top = getattr(sh, "top", 0) or 0
        height = getattr(sh, "height", 0) or 0
        return top + height / 2

    # 1. Prefer CENTER_TITLE - the centered title placeholder (main content on many slides)
    if PP_PLACEHOLDER:
        for ph in slide.placeholders:
            if _skip_placeholder(ph):
                continue
            pf = getattr(ph, "placeholder_format", None)
            if pf and getattr(pf, "type", None) == PP_PLACEHOLDER.CENTER_TITLE:
                if ph.has_text_frame and ph.text and ph.text.strip():
                    t = ph.text.strip()
                    if not _is_bad_title(t):
                        return t
                break

    # 2. Prefer shape in vertical CENTER (25%-75%) - the main content, not top header
    candidates = []
    for sh in slide.shapes:
        if _skip_placeholder(sh):
            continue
        if not (sh.has_text_frame and sh.text and sh.text.strip()):
            continue
        cy = _center_y(sh)
        if cy < slide_height * 0.2:  # skip top 20% (header/master title zone)
            continue
        if cy > slide_height * 0.85:  # skip bottom 15% (footer zone)
            continue
        t = sh.text.strip()
        if _is_bad_title(t):
            continue
        dist_from_center = abs(cy - slide_height / 2)
        candidates.append((dist_from_center, t))
    if candidates:
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]

    # 3. Fallback: TITLE or SUBTITLE placeholder (only if not bad)
    if PP_PLACEHOLDER:
        for ph in slide.placeholders:
            if _skip_placeholder(ph):
                continue
            pf = getattr(ph, "placeholder_format", None)
            if pf and getattr(pf, "type", None) in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            ):
                if ph.has_text_frame and ph.text and ph.text.strip():
                    t = ph.text.strip()
                    if not _is_bad_title(t):
                        return t
                break

    # 4. Last resort: any non-footer shape with text (skip bad titles)
    for sh in slide.shapes:
        if _skip_placeholder(sh):
            continue
        if sh.has_text_frame and sh.text and sh.text.strip():
            if _center_y(sh) < slide_height * 0.85:
                t = sh.text.strip()
                if not _is_bad_title(t):
                    return t
    return ""


def _get_section_name(prs, slide_indices: list[int], section_key: str) -> str:
    """Get section name from first substantive slide title (e.g. 'The Agile Mindset: Achieving Flow')."""
    for idx in slide_indices:
        if idx > len(prs.slides):
            continue
        title = _get_slide_title(prs.slides[idx - 1], prs)
        if not title or _is_bad_title(title):
            continue
        t = title.strip()
        if t and not re.match(r"^\d+$", t) and len(t) < 100:
            return t
    return f"Section_{section_key}"


def _detect_inner_sections(prs) -> list[tuple[str, list[int], str]]:
    """
    Scan slides for SECTION X.Y or SECTION X titles.
    Return [(section_key, [slide_indices], section_name), ...].
    section_name = first substantive slide title (e.g. 'The Agile Mindset: Achieving Flow').
    """
    groups: list[tuple[str, list[int]]] = []
    current_key: str | None = None
    current_slides: list[int] = []
    first_section_seen = False

    for i, slide in enumerate(prs.slides, 1):
        all_text = _get_all_slide_text(slide)
        m = _SECTION_RE.search(all_text)
        if m:
            key = m.group(1)
            if current_key is not None:
                groups.append((current_key, current_slides))
            elif not first_section_seen:
                current_slides = list(range(1, i)) + [i]
            else:
                current_slides = [i]
            current_key = key
            first_section_seen = True
        else:
            if current_key is not None:
                current_slides.append(i)

    if current_key is not None:
        groups.append((current_key, current_slides))

    return [(k, s, _get_section_name(prs, s, k)) for k, s in groups]


def _section_key_ge(key: str, start: str) -> bool:
    """True if key >= start (e.g. '4.1' >= '4.1', '4.2' >= '4.1')."""
    def to_tuple(s: str) -> tuple:
        return tuple(int(x) for x in s.split("."))
    try:
        return to_tuple(key) >= to_tuple(start)
    except (ValueError, AttributeError):
        return True


def split_pptx_by_section(
    src_path: Path,
    output_dir: Path | None = None,
    start_section: str | None = None,
) -> list[Path]:
    """
    Split PowerPoint by INNER section numbers (from slide titles).
    Returns list of created file paths.
    """
    if win32com is None:
        raise RuntimeError("pywin32 required. pip install pywin32")
    if Presentation is None:
        raise RuntimeError("python-pptx required. pip install python-pptx")

    src_path = Path(src_path).resolve()
    if not src_path.is_file():
        raise FileNotFoundError(f"File not found: {src_path}")

    out_dir = Path(output_dir).resolve() if output_dir else src_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    # Detect inner sections from slide titles (python-pptx)
    prs = Presentation(str(src_path))
    groups = _detect_inner_sections(prs)

    if start_section:
        groups = [(k, s, n) for k, s, n in groups if _section_key_ge(k, start_section)]
        if not groups:
            print(f"No sections >= {start_section} found.")
            return []

    if not groups:
        print("No SECTION X.Y or SECTION X titles found in slides.")
        return []

    try:
        import pythoncom
        pythoncom.CoInitialize()
    except ImportError:
        pass

    ppt = win32com.client.gencache.EnsureDispatch("PowerPoint.Application")
    ppt.Visible = 1

    created: list[Path] = []
    pres = None

    try:
        pres = ppt.Presentations.Open(str(src_path), WithWindow=False)
        time.sleep(0.3)

        for section_key, slide_indices, section_name in groups:
            if not slide_indices:
                continue

            safe_name = _sanitize_filename(section_name)
            out_path = out_dir / f"{safe_name}.pptx"
            n = 1
            while out_path in created or out_path.exists():
                n += 1
                out_path = out_dir / f"{safe_name}_{n}.pptx"

            print(f"[{len(created)+1}/{len(groups)}] Creating: {out_path.name} (slides {slide_indices[0]}-{slide_indices[-1]})")
            sys.stdout.flush()

            fd, temp_abs = tempfile.mkstemp(suffix=".pptx", prefix="split_")
            os.close(fd)
            temp_path = Path(temp_abs)
            pres.SaveCopyAs(temp_abs)
            time.sleep(0.2)

            temp_pres = None
            try:
                temp_pres = ppt.Presentations.Open(str(temp_path), WithWindow=False)
                time.sleep(0.2)

                # Delete slides NOT in this group (from end to start to preserve indices)
                total = temp_pres.Slides.Count
                to_delete = [j for j in range(1, total + 1) if j not in slide_indices]
                for j in reversed(to_delete):
                    temp_pres.Slides.Range([j]).Delete()

                temp_pres.SaveAs(str(out_path))
                created.append(out_path)
                print(f"       Created: {out_path}")
                sys.stdout.flush()
            finally:
                if temp_pres:
                    try:
                        temp_pres.Close()
                    except Exception:
                        pass
                    time.sleep(0.2)

            try:
                if temp_path.exists():
                    temp_path.unlink()
            except OSError:
                pass

        return created

    finally:
        if pres:
            try:
                pres.Close()
            except Exception:
                pass
        try:
            ppt.Quit()
        except Exception:
            pass


def main():
    args = sys.argv[1:]
    src = None
    out_dir = None
    start_section = None
    i = 0
    while i < len(args):
        if args[i] == "--output-dir" and i + 1 < len(args):
            out_dir = args[i + 1]
            i += 2
        elif args[i] == "--start-section" and i + 1 < len(args):
            start_section = args[i + 1]
            i += 2
        elif not args[i].startswith("--"):
            src = args[i]
            i += 1
        else:
            i += 1

    if not src:
        print("Usage: python split_pptx_by_section.py <path_to.pptx> [--output-dir <dir>] [--start-section 4.1]")
        sys.exit(1)

    src_path = Path(src)
    if not src_path.is_absolute():
        # Try cwd and common roots
        for base in [Path.cwd(), Path(__file__).resolve().parent.parent.parent]:
            cand = base / src
            if cand.exists():
                src_path = cand
                break

    if not src_path.exists():
        print(f"File not found: {src_path}")
        sys.exit(1)

    print(f"Source: {src_path}")
    if start_section:
        print(f"Starting from section: {start_section}")
    print("Splitting by section...")
    try:
        created = split_pptx_by_section(
            src_path,
            Path(out_dir) if out_dir else None,
            start_section=start_section,
        )
        print(f"\nCreated {len(created)} file(s).")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
